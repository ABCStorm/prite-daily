#!/usr/bin/env python3
"""
Extract PRITE questions from the PowerPoint decks into structured JSON + images.

Model (confirmed with the deck owner):
  - A QUESTION slide has a stem with A-E choices and a short answer text box.
  - The slides that FOLLOW a question (image-bearing and/or written prose) are
    its EXPLANATION. Roughly one per 2-3 questions. Several decks (2016, 2017,
    2020, 2021, 2022) contain written-out text explanations ("Why B: ...",
    "Answer: B — ...", "Explanation: ..."); those are explanation slides, NOT
    questions, and must not be mistaken for questions just because they're wordy.

A slide is treated as a REAL QUESTION only if we can recover option structure:
  (a) >= 3 lettered options (A./B./C. ... incl. concatenated "A.xB.y"), OR
  (b) a lettered answer box (e.g. "D. ...") + 3-6 short unlettered bullet lines
      whose first line is the question (ends in '?' or ':')  -- the 2025 style.
Anything else following a question becomes that question's explanation.

Output per question:
  { deck, year, q_index, slide_number, stem, options[{letter,text}],
    answer_letter, answer_text, answer_source, explanation_text,
    images[paths], flags[] }

Uncertain parses are recorded in `flags` so QA is targeted.

Usage:
  python extract.py "/path/to/deck-folder" --out output
"""
import sys
import os
import re
import json
import glob
import argparse
from pptx import Presentation
from pptx.oxml.ns import qn

LETTERS = ["A", "B", "C", "D", "E", "F", "G"]

OPT_LINE_RE = re.compile(r"^\s*([A-G])[\.\)]\s+(.*\S)\s*$")
# inline boundary: a non-space char (not itself a letter marker) then "B. "
INLINE_SPLIT_RE = re.compile(r"(?<=\S)(?=[A-G]\.\s)")
ANSWER_RE = re.compile(r"^\s*(?:answer\s*[:\-—]?\s*)?([A-G])\s*[\.\)\:\-—]\s*(.*)$", re.I)
ANSWER_BARE_RE = re.compile(r"^\s*([A-G])[\.\)]\s*$")
# A stem starting like this is a written explanation slide, not a question —
# even when it re-lists the options (common in 2020/2021).
EXPL_PREFIX_RE = re.compile(
    r"^(why\s+[a-g]\b|why\s+the\s+answer|answer\s*[:\-—]|answer\s+is\b|"
    r"answer\s+[a-g]\b|explanation\b|rationale\b|because\s+[a-g]\b)",
    re.I,
)

MAX_OPT_LEN = 170  # real option choices rarely exceed this; prose bullets do


def shape_text(shape):
    return shape.text_frame.text if shape.has_text_frame else ""


def entrance_animated_spids(slide):
    """Shape ids on this slide that have an ENTRANCE (appear-on-click) animation.
    Such shapes are hidden until the presenter clicks, so on a question slide
    they are explanation reveals — NOT part of the question stem. Parsed from the
    slide's <p:timing> tree (python-pptx doesn't model animations)."""
    spids = set()
    timing = slide.element.find(qn("p:timing"))
    if timing is None:
        return spids
    for cTn in timing.iter(qn("p:cTn")):
        if cTn.get("presetClass") == "entr":  # entrance effect
            for tgt in cTn.iter(qn("p:spTgt")):
                spid = tgt.get("spid")
                if spid:
                    spids.add(spid)
    return spids


def is_animated(shape, anim_spids):
    try:
        return str(shape.shape_id) in anim_spids
    except Exception:
        return False


def is_picture(shape):
    try:
        return shape.shape_type == 13
    except Exception:
        return False


def clean(s):
    if not s:
        return ""
    s = s.replace(" ", " ").replace("​", "").replace("﻿", "")
    s = s.replace("­", "")
    # normalize soft line breaks (vertical tab / form feed / unicode line & para
    # separators) that PowerPoint stores inside a run, so line-based parsing
    # sees option boundaries. 2021/2022 use \x0b between choices.
    s = s.replace("\x0b", "\n").replace("\x0c", "\n")
    s = s.replace(" ", "\n").replace(" ", "\n")
    s = re.sub(r"[ \t]+", " ", s)
    s = re.sub(r"\n{3,}", "\n\n", s)
    return s.strip()


def oneline(s):
    """Collapse internal newlines to single spaces (for stems / option text)."""
    return re.sub(r"\s*\n\s*", " ", s or "").strip()


def expand_sequential(opts):
    """Split options where the NEXT option's letter got glued onto the previous
    one with a period and no line break, e.g. an option text
    'Locus coeruleus C.Tuberomammillary' -> 'Locus coeruleus' + (C)'Tuberomammillary'.
    Only the next expected letter with an explicit '.'/')' is split (so 'Vitamin
    D deficiency' is left alone)."""
    out = []
    queue = [list(o) for o in opts]
    i = 0
    while i < len(queue):
        letter, text = queue[i]
        nxt = chr(ord(letter) + 1)
        m = re.search(rf"\s+{nxt}[\.\)](?=\s*\S)", text) if nxt <= "H" else None
        if m:
            head = clean(text[: m.start()])
            tail = clean(text[m.end():])
            if head and tail:
                out.append((letter, head))
                queue.insert(i + 1, [nxt, tail])
                i += 1
                continue
        out.append((letter, clean(text)))
        i += 1
    return out


def lettered_options(text):
    """Find lettered options via line-based and inline strategies. Returns list
    of (letter, text) or [] if it doesn't look like a real option set."""
    text = text.replace("\r", "\n")
    lines = text.split("\n")

    opts = []
    first = None
    for i, ln in enumerate(lines):
        m = OPT_LINE_RE.match(ln)
        if m:
            if first is None:
                first = i
            opts.append([m.group(1), clean(m.group(2))])
        elif first is not None and ln.strip() and opts:
            # wrapped continuation of previous option
            opts[-1][1] = clean(opts[-1][1] + " " + ln.strip())

    if len(opts) >= 3:
        stem = clean("\n".join(lines[:first]))
        # recover a leading option that got glued onto the stem tail, e.g.
        # "...the exome? A 2%" when B/C/D/E are on their own lines.
        if opts[0][0] != "A":
            mlead = re.search(r"(?:^|[\?\:\s])A[\.\)]?\s+(\S.+)$", stem)
            if mlead and LETTERS.index(opts[0][0]) == 1:  # options start at B
                stem = clean(stem[: mlead.start()])
                opts.insert(0, ["A", clean(mlead.group(1))])
        opts = expand_sequential(opts)
        return stem, [(l, t) for l, t in opts]

    # inline concatenated ("...A. xB. yC. z")
    joined = " ".join(lines)
    m = re.search(r"(?<![A-Za-z])([A-G])[\.\)]\s", joined)
    if m:
        stem = clean(joined[: m.start()])
        rest = joined[m.start():]
        parts = [p for p in re.split(INLINE_SPLIT_RE, rest) if p.strip()]
        io = []
        for p in parts:
            mm = OPT_LINE_RE.match(p.strip()) or re.match(r"^([A-G])[\.\)]\s*(.*)$", p.strip())
            if mm:
                io.append((mm.group(1), clean(mm.group(2))))
        if len(io) >= 3:
            return stem, io

    return None


def unlettered_options(text, answer_letter, has_answer_box):
    """2025 style: a question line then N short bullet lines with no A-E letters.
    The answer box may also be unlettered (just the winning choice's text), so we
    only require that SOME answer box exists; structure guards against eating
    explanation prose: first line ends in '?'/':' and 3-6 short bullets follow."""
    if not has_answer_box:
        return None
    lines = [clean(l) for l in text.replace("\r", "\n").split("\n") if clean(l)]
    if len(lines) < 4:
        return None
    # question line: first line ending with ? or :
    q_end = None
    for i, ln in enumerate(lines):
        if ln.endswith("?") or ln.endswith(":"):
            q_end = i
            break
    if q_end is None:
        return None
    stem = clean(" ".join(lines[: q_end + 1]))
    bullets = lines[q_end + 1:]
    if not (3 <= len(bullets) <= 6):
        return None
    if any(len(b) > MAX_OPT_LEN for b in bullets):
        return None
    if answer_letter in LETTERS and len(bullets) < LETTERS.index(answer_letter) + 1:
        return None
    return stem, [(LETTERS[i], b) for i, b in enumerate(bullets)]


ANS_PREFIX_RE = re.compile(r"^\s*answer\s*[:\-—]?\s*", re.I)
# a single leading option letter followed by a word boundary (so "Cerebral" is
# NOT read as letter C, but "C. Anterior", "C anterior", and bare "C" are).
SINGLE_LETTER_RE = re.compile(r"^([A-H])\b[\.\)\:\-—\s]*(.*)$")
# a list of letters: "A, B, C", "A/B", "B and D"
LETTER_LIST_RE = re.compile(r"^[A-H]([\s,/&]+(and\s+)?[A-H])+$", re.I)
# concatenated letters with no separators: "ACG"
LETTER_RUN_RE = re.compile(r"^[A-H]{2,8}$")


def answer_core(raw):
    return ANS_PREFIX_RE.sub("", oneline(clean(raw))).strip().rstrip(".").strip()


def parse_letter_list(core):
    """Return a list of letters if `core` is purely a multi-answer key, else None."""
    if LETTER_LIST_RE.match(core):
        seen = []
        for l in re.findall(r"[A-H]", core, re.I):
            L = l.upper()
            if L not in seen:
                seen.append(L)
        return seen
    if LETTER_RUN_RE.match(core):
        return list(core)
    return None


def first_letter(raw):
    """Best-guess single leading letter of an answer box (for sizing options)."""
    core = answer_core(raw)
    m = SINGLE_LETTER_RE.match(core)
    if m:
        return m.group(1).upper()
    ll = parse_letter_list(core)
    return ll[0] if ll else None


def looks_like_answer(t):
    core = answer_core(t)
    return bool(SINGLE_LETTER_RE.match(core) or parse_letter_list(core))


def find_answer_box(texts, stem_text):
    """Among text shapes (excluding the stem), pick the answer-key text:
    prefer a short answer-looking candidate, else the shortest."""
    cands = [t for t in texts if t is not stem_text]
    if not cands:
        return ""
    cands.sort(key=lambda t: (0 if (looks_like_answer(t) and len(t) <= MAX_OPT_LEN) else 1, len(t)))
    return cands[0]


def parse_question(texts):
    """Try to interpret a slide's texts as a real question.
    Returns dict(stem, options, answer_letter, answer_text, answer_raw,
    answer_source, flags) or None."""
    if not texts:
        return None
    stem_text = max(texts, key=len)
    answer_raw = find_answer_box(texts, stem_text)
    ans_letter = first_letter(answer_raw)

    flags = []
    has_ans = bool(answer_raw.strip())
    # Some decks (esp. 2025) split stem and options into SEPARATE text boxes.
    # Try the longest box first; if that yields no options, retry on the
    # combined non-answer boxes (in length order, stem first).
    rest_boxes = [t for t in texts if t is not stem_text and t != answer_raw]
    combined = "\n".join([stem_text] + sorted(rest_boxes, key=len, reverse=True))

    res = lettered_options(stem_text)
    source = "lettered"
    if not res and rest_boxes:
        res = lettered_options(combined)
    if not res:
        res = unlettered_options(stem_text, ans_letter, has_ans)
        if not res and rest_boxes:
            res = unlettered_options(combined, ans_letter, has_ans)
        source = "unlettered"
        if res:
            flags.append("options_unlettered")
    if not res:
        return None

    stem, options = res
    if not stem:
        stem = stem_text
    stem = oneline(stem)
    if EXPL_PREFIX_RE.match(stem):
        return None  # written explanation that happens to re-list options
    stem = re.sub(r"^\d{1,3}[\.\)]\s+", "", stem)  # strip leading question number
    options = [(l, oneline(t)) for l, t in options]

    # resolve answer
    ans = resolve_answer(answer_raw, options)
    flags += ans["flags"]
    if len(options) < 4:
        flags.append(f"only_{len(options)}_options")
    if len(options) > 6:
        flags.append(f"many_{len(options)}_options")

    return {
        "stem": stem,
        "options": [{"letter": l, "text": t} for l, t in options],
        "answer_letter": ans["answer_letter"],
        "answer_letters": ans["answer_letters"],
        "multi_select": ans["multi_select"],
        "answer_text": oneline(ans["answer_text"]),
        "answer_raw": oneline(clean(answer_raw)),
        "answer_source": ans["answer_source"],
        "flags": flags,
    }


def _ans(letter, letters, text, multi, source, flags):
    return {
        "answer_letter": letter,
        "answer_letters": letters,
        "answer_text": text,
        "multi_select": multi,
        "answer_source": source,
        "flags": flags,
    }


def resolve_answer(answer_raw, options):
    core = answer_core(answer_raw)
    valid = {l for l, _ in options}
    if not core:
        return _ans(None, [], "", False, "none", ["no_answer_box"])

    # multi-select ("A, B, C", "ACG", "B and D")
    letters = parse_letter_list(core)
    if letters and len(letters) >= 2:
        keep = [l for l in letters if l in valid] or letters
        txt = "; ".join((f"{l}. {opt_text(options, l)}").strip() for l in keep)
        return _ans(keep[0], keep, oneline(txt), True, "multi", ["multi_select"])

    # single leading letter (validated against the option set)
    m = SINGLE_LETTER_RE.match(core)
    if m and m.group(1).upper() in valid:
        L = m.group(1).upper()
        rest = oneline(m.group(2))
        return _ans(L, [L], rest or opt_text(options, L), False, "letter", [])

    # an embedded "X. <option text>" where X is valid and its text matches an
    # option (e.g. answer box 'lnsula D. Medial pre-frontal cortex' -> D).
    for em in re.finditer(r"\b([A-H])[\.\)]\s*(\S[^A-H]*)", core):
        L = em.group(1).upper()
        et = normtext(em.group(2))
        if L in valid and len(et) >= 6 and et in normtext(opt_text(options, L)):
            return _ans(L, [L], opt_text(options, L), False, "letter-embedded", ["answer_embedded_letter"])

    # no letter — match the answer text against the option texts
    norm = normtext(core)
    for l, t in options:
        if normtext(t) == norm:
            return _ans(l, [l], clean(t), False, "text-match", [])
    for l, t in options:
        if norm and (normtext(t) in norm or norm in normtext(t)):
            return _ans(l, [l], clean(t), False, "text-fuzzy", ["answer_fuzzy_match"])
    return _ans(None, [], core, False, "unmatched", ["answer_unmatched"])


def opt_text(options, letter):
    for l, t in options:
        if l == letter:
            return clean(t)
    return ""


def normtext(s):
    return re.sub(r"[^a-z0-9]", "", (s or "").lower())


def gather(slide):
    texts, images = [], []
    for shape in slide.shapes:
        if is_picture(shape):
            images.append(shape)
        t = clean(shape_text(shape))
        if t:
            texts.append(t)
    return texts, images


def is_title(texts, images):
    if texts and not images and max((len(t) for t in texts), default=0) < 40:
        j = " ".join(texts).lower()
        if "prite" in j or j.strip().lower().startswith("part"):
            return True
    return False


def export_images(images, out_img_dir, prefix):
    paths = []
    for i, shape in enumerate(images):
        try:
            img = shape.image
            ext = img.ext or "png"
            fname = f"{prefix}_{i+1}.{ext}"
            with open(os.path.join(out_img_dir, fname), "wb") as f:
                f.write(img.blob)
            paths.append(f"images/{os.path.basename(out_img_dir)}/{fname}")
        except Exception as e:
            paths.append(f"<image-export-failed: {e}>")
    return paths


def extract_deck(path, out_dir):
    deck = os.path.basename(path)
    ym = re.search(r"(20\d\d)", deck)
    year = ym.group(1) if ym else deck
    prs = Presentation(path)
    img_dir = os.path.join(out_dir, "images", year)
    os.makedirs(img_dir, exist_ok=True)

    questions = []
    current = None
    q_index = 0

    for sidx, slide in enumerate(prs.slides, start=1):
        texts, images = gather(slide)
        if not texts and not images:
            continue
        if is_title(texts, images):
            continue

        parsed = parse_question(texts)
        anim = entrance_animated_spids(slide)

        if parsed:
            if current:
                questions.append(current)
            q_index += 1
            # On a question slide, an image is a genuine question FIGURE (CT scan,
            # pedigree, brain section) only if it is visible immediately — i.e. it
            # has NO entrance animation. Images that "appear on click" are
            # explanation reveals (e.g. a pasted answer screenshot) and must stay
            # hidden until the user answers, so they go to explanation_images.
            visible = [im for im in images if not is_animated(im, anim)]
            reveal = [im for im in images if is_animated(im, anim)]
            figure_images = export_images(visible, img_dir, f"{year}_q{q_index:03d}_fig") if visible else []
            onslide_expl = export_images(reveal, img_dir, f"{year}_q{q_index:03d}_expA") if reveal else []
            flags = parsed["flags"]
            if figure_images:
                flags = flags + ["has_figure"]
            if onslide_expl:
                flags = flags + ["onslide_reveal"]
            current = {
                "deck": deck,
                "year": year,
                "q_index": q_index,
                "slide_number": sidx,
                "stem": parsed["stem"],
                "options": parsed["options"],
                "answer_letter": parsed["answer_letter"],
                "answer_letters": parsed["answer_letters"],
                "multi_select": parsed["multi_select"],
                "answer_text": parsed["answer_text"],
                "answer_source": parsed["answer_source"],
                "answer_raw": parsed["answer_raw"],
                "explanation_text": "",
                "figure_images": figure_images,
                "explanation_images": list(onslide_expl),
                "flags": flags,
            }
        else:
            # explanation material on the slides FOLLOWING a question
            if current is None:
                continue
            prefix = f"{year}_q{current['q_index']:03d}_exp{sidx}"
            if images:
                current["explanation_images"] += export_images(images, img_dir, prefix)
            exp = " ".join(texts).strip()
            if exp:
                current["explanation_text"] = clean(
                    (current["explanation_text"] + "\n" + exp).strip()
                )

    if current:
        questions.append(current)
    return year, questions


def load_corrections():
    """Manual overrides keyed by (year, slide_number), applied after extraction
    so hand-fixes survive a re-run. See corrections.json."""
    path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "corrections.json")
    if not os.path.exists(path):
        return {}
    data = json.load(open(path))
    return {(str(c["year"]), int(c["slide_number"])): c for c in data}


def apply_corrections(qs, corrections):
    n = 0
    for q in qs:
        c = corrections.get((str(q["year"]), int(q["slide_number"])))
        if not c:
            continue
        q.update(c.get("set", {}))
        # drop now-stale error flags and mark as hand-corrected
        clean_flags = [f for f in q.get("flags", [])
                       if f not in ("answer_unmatched", "no_answer_box",
                                    "answer_fuzzy_match", "only_3_options", "only_2_options")]
        if "manual_corrected" not in clean_flags:
            clean_flags.append("manual_corrected")
        q["flags"] = clean_flags
        n += 1
    return n


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("folder")
    ap.add_argument("--out", default="output")
    args = ap.parse_args()

    decks = sorted(
        d for d in glob.glob(os.path.join(args.folder, "*.pptx"))
        if not os.path.basename(d).startswith("~$")  # skip Office lock/temp files
    )
    os.makedirs(args.out, exist_ok=True)
    corrections = load_corrections()

    all_q = []
    summary = []
    applied = 0
    for d in decks:
        year, qs = extract_deck(d, args.out)
        applied += apply_corrections(qs, corrections)
        with open(os.path.join(args.out, f"questions_{year}.json"), "w") as f:
            json.dump(qs, f, indent=2, ensure_ascii=False)
        flagged = sum(1 for q in qs if q["flags"])
        no_ans = sum(1 for q in qs if not q["answer_letter"])
        multi = sum(1 for q in qs if q["multi_select"])
        fig = sum(1 for q in qs if q["figure_images"])
        with_exp = sum(1 for q in qs if q["explanation_text"] or q["explanation_images"])
        summary.append((os.path.basename(d), len(qs), flagged, no_ans, multi, fig, with_exp))
        all_q.extend(qs)

    with open(os.path.join(args.out, "questions_all.json"), "w") as f:
        json.dump(all_q, f, indent=2, ensure_ascii=False)

    print(f"{'deck':<20}{'Qs':>5}{'flag':>6}{'no-ans':>7}{'multi':>6}{'figs':>6}{'w/expl':>7}")
    print("-" * 62)
    tot = [0, 0, 0, 0, 0, 0]
    for name, nq, fl, na, mu, fg, we in summary:
        print(f"{name:<20}{nq:>5}{fl:>6}{na:>7}{mu:>6}{fg:>6}{we:>7}")
        tot = [tot[0]+nq, tot[1]+fl, tot[2]+na, tot[3]+mu, tot[4]+fg, tot[5]+we]
    print("-" * 62)
    print(f"{'TOTAL':<20}{tot[0]:>5}{tot[1]:>6}{tot[2]:>7}{tot[3]:>6}{tot[4]:>6}{tot[5]:>7}")
    print(f"\nApplied {applied} manual correction(s) from corrections.json")
    print(f"Wrote {len(all_q)} questions to {args.out}/questions_all.json")


if __name__ == "__main__":
    main()
