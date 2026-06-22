#!/usr/bin/env python3
"""
Inspect PRITE PowerPoint decks to understand their structure BEFORE writing
the real extractor. Prints, per deck: slide count, how many shapes/text frames
per slide, image counts, and a sample of the actual text from the first several
slides so we can see how a "question" is laid out and how consistent it is.

Usage:
    python inspect_decks.py "/path/to/folder/with/decks"
"""
import sys
import os
import glob
from collections import Counter
from pptx import Presentation
from pptx.util import Emu


def shape_summary(shape):
    kind = shape.shape_type
    has_text = shape.has_text_frame and shape.text_frame.text.strip()
    is_pic = (str(kind) == "PICTURE (13)") or shape.shape_type == 13
    return kind, bool(has_text), is_pic


def deck_overview(path):
    name = os.path.basename(path)
    try:
        prs = Presentation(path)
    except Exception as e:
        print(f"\n### {name}\n  !! FAILED TO OPEN: {e}")
        return

    slides = prs.slides
    n = len(slides)
    text_shape_counts = []
    pic_counts = []
    total_pics = 0
    chars_per_slide = []

    for slide in slides:
        tshapes = 0
        pics = 0
        chars = 0
        for shape in slide.shapes:
            _, has_text, is_pic = shape_summary(shape)
            if has_text:
                tshapes += 1
                chars += len(shape.text_frame.text)
            if is_pic:
                pics += 1
        text_shape_counts.append(tshapes)
        pic_counts.append(pics)
        total_pics += pics
        chars_per_slide.append(chars)

    print(f"\n### {name}")
    print(f"  slides: {n}")
    print(f"  slide w x h (in): {Emu(prs.slide_width).inches:.1f} x {Emu(prs.slide_height).inches:.1f}")
    print(f"  text-shapes/slide: dist {dict(Counter(text_shape_counts).most_common(6))}")
    print(f"  images/slide:      dist {dict(Counter(pic_counts).most_common(6))}  (total {total_pics})")
    if chars_per_slide:
        avg = sum(chars_per_slide) / len(chars_per_slide)
        print(f"  avg chars/slide: {avg:.0f}  (min {min(chars_per_slide)}, max {max(chars_per_slide)})")

    # show the raw text of the first 4 slides so we can see the layout pattern
    print("  --- first 4 slides, raw text per shape ---")
    for i, slide in enumerate(slides):
        if i >= 4:
            break
        print(f"  [slide {i+1}]")
        for j, shape in enumerate(slide.shapes):
            if shape.has_text_frame and shape.text_frame.text.strip():
                txt = shape.text_frame.text.strip().replace("\n", " / ")
                if len(txt) > 240:
                    txt = txt[:240] + "…"
                print(f"     shape{j} [{shape.shape_type}]: {txt}")
            elif (str(shape.shape_type) == "PICTURE (13)") or shape.shape_type == 13:
                print(f"     shape{j} [PICTURE]")


def main():
    folder = sys.argv[1] if len(sys.argv) > 1 else "."
    decks = sorted(glob.glob(os.path.join(folder, "*.pptx")))
    if not decks:
        print(f"No .pptx files found in {folder}")
        sys.exit(1)
    print(f"Found {len(decks)} decks in {folder}")
    for d in decks:
        deck_overview(d)


if __name__ == "__main__":
    main()
